import os
import csv
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import mokepy as mkpy  # Assuming MOKEAnalyzer is in a separate file
import numpy as np
from importlib import reload
from scipy.stats import norm

import scienceplots
from matplotlib.colors import Normalize
from matplotlib.cm import ScalarMappable
from typing import List, Dict, Optional
import natsort

plt.style.use(['science', 'ieee', 'no-latex'])

reload(mkpy)


class BatchAnalyzer:
    def __init__(self, folder_path):
        self.folder_path = folder_path
        self.analyzers = {}
        self.summary_data = []

    def process_files(self, field_threshold=1500):
        """Process all .csv files in the specified folder."""
        filenames = os.listdir(self.folder_path)
        filenames_s = natsort.natsorted(filenames)
        for filename in filenames_s:
            if filename.endswith('.CSV') and 'coercive_fields_summary' not in filename:
                print(filename)
                file_path = os.path.join(self.folder_path, filename)
                analyzer = mkpy.MOKEAnalyzer(file_path)
                analyzer.read_data()
                analyzer.fit_background(field_threshold=field_threshold)
                analyzer.normalize_data()
                analyzer.find_coercive_field()
                analyzer.fit_switching_field_distribution()
                analyzer.calc_Mr(field_threshold=field_threshold)
                analyzer.calc_squareness(field_threshold=field_threshold)
                
                self.analyzers[filename] = analyzer
                sweep_rate = self.extract_sweep_rate(filename)
                
                # Collect summary data
                self.summary_data.append({
                    'Filename': filename,
                    'Hc1': analyzer.Hc_summary['Hc1'],
                    'Hc2': analyzer.Hc_summary['Hc2'],
                    'Hc_mean': analyzer.Hc_summary['Hc_mean'],
                    'sfd_amp1':analyzer.sfd_fit_params['popt_pos'][0],
                    'sfd_amp2':analyzer.sfd_fit_params['popt_neg'][0],
                    'sfd_mean1':analyzer.sfd_fit_params['popt_pos'][1],
                    'sfd_mean2':analyzer.sfd_fit_params['popt_neg'][1],
                    'sfd_std1':analyzer.sfd_fit_params['popt_pos'][2],
                    'sfd_std2':analyzer.sfd_fit_params['popt_neg'][2],
                    'sfd_std_mean':(abs(analyzer.sfd_fit_params['popt_neg'][2])+abs(analyzer.sfd_fit_params['popt_pos'][2]))/2,
                    'sfd_const1':analyzer.sfd_fit_params['popt_pos'][3],
                    'sfd_const2':analyzer.sfd_fit_params['popt_neg'][3],
                    'Mr_mean':analyzer.Mr_summary['Mr_mean'],
                    'Mr_std':analyzer.Mr_summary['Mr_std'],
                    'MrMs':analyzer.Mr_summary['MrMs'],
                    'Ms_mean':analyzer.Ms_summary['Ms_mean'],
                    'Ms_std':analyzer.Ms_summary['Ms_std'],
                    'sweep_rate':sweep_rate
                })

    def extract_sweep_rate(self, filename):
        """Extract the field sweep rate from the filename."""
        parts = filename.split('-')
        field_part = next(part for part in parts if 'kOe' in part)
        time_part = next(part for part in parts if 's' in part and 'kOe' not in part)
        
        total_field = float(field_part.replace('kOe', '')) * 1000  # Convert to Oe
        sweep_time = float(time_part.replace('s.CSV', ''))
        
        sweep_rate = (4 * total_field) / sweep_time  # Multiply by 4 for full sweep
        # Alternative calculation
        field = self.analyzers[filename].data['H']
        field = np.absolute(np.diff(field))
        total_field = np.sum(field)
        sweep_rate = total_field/sweep_time
        return sweep_rate
    
    def plot_parameter_vs_sweep_rate(self, parameter, log_scale=False):
        """Plot a parameter against sweep rate."""
        sweep_rates = []
        for i,data in enumerate(self.summary_data):
            sweep_rate = self.extract_sweep_rate(data['Filename'])
            sweep_rates.append(sweep_rate)
            #data['sweep_rate'] = sweep_rate
            #self.summary_data[i]['sweep_rate'] = sweep_rate
        #sweep_rates = [self.extract_sweep_rate(data['Filename']) for data in self.summary_data]
        parameter_values = [data[parameter] for data in self.summary_data]
        
        plt.figure(figsize=(4, 3))
        if log_scale:
            plt.plot(np.log(sweep_rates), parameter_values, 'o')
            plt.xlabel('Sweep Rate (Oe/s) - Log Scale')
        else:
            plt.plot(sweep_rates, parameter_values, 'o')
            plt.xlabel('Sweep Rate (Oe/s)')
        
        plt.ylabel(parameter)
        plt.title(f'{parameter} vs Sweep Rate')
        plt.grid(True)
        
        plot_filename = f'{parameter.lower().replace(" ", "_")}_vs_sweep_rate{"_log" if log_scale else ""}.png'
        plt.savefig(os.path.join(self.folder_path, plot_filename))
        plt.close()
        return(plot_filename)

    def plot_all_parameters_vs_sweep_rate(self):
        parameters = ['Hc_mean', 'sfd_std_mean', 'Mr_mean', 'MrMs']
        plot_filenames = []
        for param in parameters:
            plot_filenames.append(self.plot_parameter_vs_sweep_rate(param, log_scale=False))
            plot_filenames.append(self.plot_parameter_vs_sweep_rate(param, log_scale=True))
        return plot_filenames

    def calculate_summary_statistics(self):
        hc_mean_values = [data['Hc_mean'] for data in self.summary_data]
        avg_hc = np.mean(hc_mean_values)
        
        sfd_width_values = [(data['sfd_std1'] + data['sfd_std2']) / 2 for data in self.summary_data]
        avg_sfd_width = np.mean(sfd_width_values)
        
        return avg_hc, avg_sfd_width

    def extract_parameter(self, filename, param_name):
        """Extracts a specified parameter from the filename or data."""
        # This is a placeholder implementation. You may need to adjust it based on your actual data structure.
        # For demonstration purposes, let's assume we can extract parameters from the filename.
        if param_name == 'Area':
            return float(filename.split('-')[3])  # Example extraction logic (adjust as needed)
        elif param_name == 'Thickness':
            return float(filename.split('-')[4])  # Example extraction logic (adjust as needed)
        else:
            return None

    def plot_combined_hysteresis_loops(self):
        """Plot hysteresis loops from all files on a single graph."""
        plt.figure(figsize=(6, 4))
        # Create a colormap
        cmap = plt.get_cmap('viridis')
        norm_c = Normalize(vmin=0, vmax=len(self.analyzers)-1)
        
        for i, (filename, analyzer) in enumerate(self.analyzers.items()):
            color = cmap(norm_c(i))
            plt.plot(analyzer.data['H'], analyzer.data['M_norm'], color=color, label=filename)
        plt.xlabel('Field (Oe)')
        plt.ylabel('Normalized Magnetization')
        plt.title('Combined Hysteresis Loops')
        plt.legend()
        plt.savefig(os.path.join(self.folder_path, 'combined_hysteresis_loops.png'))
        plt.close()

    def plot_combined_switching_field_distributions(self):
        """Plot switching field distributions from all files on a single graph."""
        plt.figure(figsize=(6, 4))

        # Create a colormap
        cmap = plt.get_cmap('viridis')
        norm_scale = Normalize(vmin=0, vmax=len(self.analyzers)-1)
        
        for i, (filename, analyzer) in enumerate(self.analyzers.items()):
            color = cmap(norm_scale(i))
            field, magnetization = analyzer.data['H'], analyzer.data['M_norm']
            mag_diff = np.diff(magnetization)
            mag_cdf = norm.cdf(mag_diff)

            sfd = np.abs(mag_diff)
            plt.plot(field[:-1], sfd, color=color,label=filename)
        plt.xlabel('Field (Oe)')
        plt.ylabel('Switching field distribution ($Oe^{-1}$)')
        plt.xlim(-1500,1500)
        plt.title('Combined Switching Field Distributions')
        plt.legend()
        plt.savefig(os.path.join(self.folder_path, 'combined_switching_field_distributions.png'))
        plt.close()

    def save_coercive_fields(self):
        """Save coercive fields to a CSV file."""
        csv_path = os.path.join(self.folder_path, 'coercive_fields_summary.csv')
        with open(csv_path, 'w', newline='') as csvfile:
            fieldnames = ['Filename', 'Hc1', 'Hc2', 'Hc_mean', 'sfd_amp1', 'sfd_amp2', 'sfd_mean1', 'sfd_mean2', 'sfd_std1', 'sfd_std2','sfd_std_mean', 'sfd_const1', 'sfd_const2','Mr_mean', 'Mr_std', 'MrMs', 'Ms_mean', 'Ms_std', 'sweep_rate']
            writer = csv.DictWriter(csvfile, fieldnames=fieldnames)
            writer.writeheader()
            for data in self.summary_data:
                writer.writerow(data)

    def create_summary_presentation(self):
        """Create a PowerPoint presentation summarizing the analysis."""
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "MOKE Analysis Summary"
        subtitle.text = f"Folder: {self.folder_path}"

        # Combined Hysteresis Loops slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Combined Hysteresis Loops"
        img_path = os.path.join(self.folder_path, 'combined_hysteresis_loops.png')
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(5))

        # Combined Switching Field Distributions slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Combined Switching Field Distributions"
        img_path = os.path.join(self.folder_path, 'combined_switching_field_distributions.png')
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(5))

        # Coercive Fields Summary slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Coercive Fields Summary"
        
        # Add a table with coercive field datf'filea
        rows = len(self.summary_data) + 1
        cols = 6  # Updated to include Area and Thickness
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(0.8 * rows)

        # Set column headers
        headers = ['Filename', 'Hc1', 'Hc2', 'Hc_mean', 'sfd_amp1', 'sfd_amp2', 'sfd_mean1', 'sfd_mean2', 'sfd_std1', 'sfd_std2', 'sfd_std_mean', 'sfd_const1', 'sfd_const2','Mr_mean', 'Mr_std', 'MrMs', 'Ms_mean', 'Ms_std', 'sweep_rate']
        
        # Update the number of columns
        cols = len(headers)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        for i, header in enumerate(headers):
            table.cell(0, i).text = header

        # Fill in the data
        for i, data in enumerate(self.summary_data, start=1):
            table.cell(i, 0).text = data['Filename']
            table.cell(i, 1).text = f"{data['Hc1']:.2f}"
            table.cell(i, 2).text = f"{data['Hc2']:.2f}"
            table.cell(i, 3).text = f"{data['Hc_mean']:.2f}"
            table.cell(i, 4).text = f"{data['sfd_amp1']:.4f}"
            table.cell(i, 5).text = f"{data['sfd_amp2']:.4f}"
            table.cell(i, 6).text = f"{data['sfd_mean1']:.4f}"
            table.cell(i, 7).text = f"{data['sfd_mean2']:.4f}"
            table.cell(i, 8).text = f"{data['sfd_std1']:.4f}"
            table.cell(i, 9).text = f"{data['sfd_std2']:.4f}"
            table.cell(i, 10).text = f"{data['sfd_const1']:.4f}"
            table.cell(i, 11).text = f"{data['sfd_const2']:.4f}"
        # Save the presentation

        plot_filenames = self.plot_all_parameters_vs_sweep_rate()
        for i in np.arange(0, len(plot_filenames),2):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = plot_filenames[i].replace('_', ' ').replace('.png', '').title()
            img_path = os.path.join(self.folder_path, plot_filenames[i])
            slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(5))
            img_path = os.path.join(self.folder_path, plot_filenames[i+1])
            slide.shapes.add_picture(img_path, Inches(6), Inches(1.5), width=Inches(5))

        # Summary slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Analysis Summary"

        # Add all plots to the slide
        plots = ['combined_hysteresis_loops.png', 'combined_switching_field_distributions.png']
        plots.extend([f for f in os.listdir(self.folder_path) if f.startswith('parameter_vs_sweep_rate')])

        left = top = Inches(1)
        width = height = Inches(3)
        for i, plot in enumerate(plots):
            img_path = os.path.join(self.folder_path, plot)
            left = Inches(1 + (i % 2) * 3.5)
            top = Inches(1.5 + (i // 2) * 3)
            slide.shapes.add_picture(img_path, left, top, width=width, height=height)

        # Add summary text
        avg_hc, avg_sfd_width = self.calculate_summary_statistics()
        text_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = f"Average Coercive Field: {avg_hc:.2f} Oe\n"
        p.text += f"Average Width of Switching Field Distribution: {avg_sfd_width:.2f} Oe"

        prs.save(os.path.join(self.folder_path, 'MOKE_Analysis_Summary.pptx'))
    
    

    def plot_coercive_fields_vs_parameter(self, parameter=[]):
        """Plots coercive fields against a specified parameter."""
        Hc_mean_val = [data['Hc_mean'] for data in self.summary_data]
        if len(parameter) != Hc_mean_val:
            raise ValueError(f"Invalid parameter '{parameter}'.Not the same length as the coercive fields.")

        x_values = [data[parameter] for data in self.summary_data]
        

        plt.figure(figsize=(5, 3))
        plt.scatter(x_values, Hc_mean_val, label='Hc1', color='k', marker='o')

        plt.xlabel()
        plt.ylabel("Coercive Field (Oe)")
        plt.title(f"Coercive Fields")
        plt.legend()
        plt.grid(True)

        # Save plot to file
        plt.savefig(os.path.join(self.folder_path, f'coercive_fields_vs_.png'))
        plt.close()



    def run_batch_analysis(self):
       """Run the complete batch analysis process."""
       self.process_files()
       self.plot_combined_hysteresis_loops()
       self.plot_combined_switching_field_distributions()
       self.save_coercive_fields()
       plot_filenames = self.plot_all_parameters_vs_sweep_rate()
       self.create_summary_presentation()

# Usage example:
# batch_analyzer = BatchAnalyzer('/path/to/your/data/folder')
# batch_analyzer.run_batch_analysis()
# batch_analyzer.plot_coercive_fields_vs_parameter(parameter='Area')  # Plot coercive fields vs Area or Thickness